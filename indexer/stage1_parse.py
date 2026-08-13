"""Stage 1 — Parse: convert a document into elements + pic_info using Docling."""

import os
import base64
import csv
import json
import logging
import mimetypes
import sys
import tempfile
from pathlib import Path
from typing import Any

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    handlers=[logging.StreamHandler(sys.stdout)],
)
log = logging.getLogger(__name__)

# Optional heavy dependencies (graceful fallback when absent)
try:
    import pandas as pd
    _PANDAS_AVAILABLE = True
except ImportError:
    _PANDAS_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Pt
    _PPTX_AVAILABLE = True
except ImportError:
    _PPTX_AVAILABLE = False

# Docling-parseable extensions handled via convert_file
_DOCLING_EXTENSIONS = {".pdf", ".docx", ".doc", ".odt", ".rtf", ".html", ".htm"}
# Extensions handled by local parsers
_EXCEL_EXTENSIONS = {".xlsx", ".xls"}
_CSV_EXTENSIONS = {".csv"}
_PPTX_EXTENSIONS = {".pptx", ".ppt"}
_JSON_EXTENSIONS = {".json"}
_IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp"}


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------

def parse(path: str, tmp_dir: str, docling) -> tuple[list[dict], list[dict]]:
    """
    Convert *path* into a flat list of elements and a list of pic_info dicts.

    Parameters
    ----------
    path     : local file path
    tmp_dir  : writable temporary directory for picture images
    docling  : DoclingInference instance (from core.docling)

    Returns
    -------
    elements : list of element dicts (type, text, page, section, …)
    pic_info : list of picture metadata dicts
    """
    log.info(f"Parsing document: {path}")
    suffix = Path(path).suffix.lower()

    if suffix in _DOCLING_EXTENSIONS:
        return _parse_via_docling(path, tmp_dir, docling)
    if suffix in _EXCEL_EXTENSIONS:
        return _parse_excel(path)
    if suffix in _CSV_EXTENSIONS:
        return _parse_csv(path)
    if suffix in _PPTX_EXTENSIONS:
        return _parse_pptx(path, tmp_dir)
    if suffix in _JSON_EXTENSIONS:
        return _parse_json(path)
    if suffix in _IMAGE_EXTENSIONS:
        return _parse_image(path, tmp_dir)
    # Fallback: treat as plain text
    return _parse_plaintext(path)


# ---------------------------------------------------------------------------
# Docling path
# ---------------------------------------------------------------------------

def _parse_via_docling(path: str, tmp_dir: str, docling) -> tuple[list, list]:
    log.info(f"Parsing via Docling: {path}")
    result = docling.convert_file(path, to_formats=["md", "json"])

    # Normalise the response envelope (docling-serve returns {"documents": [...]} or {"document": ...})
    documents = result.get("documents", result.get("document"))
    if isinstance(documents, list):
        doc = documents[0] if documents else {}
    elif isinstance(documents, dict):
        doc = documents
    else:
        doc = result  # unexpected shape — try root directly

    json_content = doc.get("json_content") or doc.get("content", {})

    if json_content:
        return _parse_docling_json(json_content, tmp_dir)

    # Fallback: parse the markdown text when JSON is unavailable
    md_content = doc.get("md_content") or doc.get("text", "")
    return _parse_markdown_text(md_content), []


def _parse_docling_json(doc_json: dict, tmp_dir: str) -> tuple[list, list]:
    """Parse a DoclingDocument JSON into elements + pic_info."""
    log.info("Parsing Docling JSON content")
    elements: list[dict] = []
    pic_info: list[dict] = []

    # Build ref → item map
    ref_map: dict[str, dict] = {}
    for key in ("texts", "tables", "pictures"):
        for item in doc_json.get(key, []):
            ref = item.get("self_ref", "")
            if ref:
                ref_map[ref] = item

    current_section = ""
    body_children = doc_json.get("body", {}).get("children", [])

    for child in body_children:
        ref = child.get("$ref", "")
        item = ref_map.get(ref)
        if item is None:
            continue

        label = item.get("label", "text")
        prov = item.get("prov") or [{}]
        page = prov[0].get("page_no", 0) if prov else 0
        text = item.get("text", "").strip()

        if label in ("section_header", "title", "page_header"):
            current_section = text
            elements.append({"type": "heading", "text": text, "page": page, "section": current_section})

        elif label == "table":
            table_data = item.get("data", {})
            elements.append({
                "type": "table",
                "text": "",
                "page": page,
                "section": current_section,
                "table_data": table_data,
                "rows": table_data.get("num_rows", 0),
                "cols": table_data.get("num_cols", 0),
            })

        elif label == "picture":
            element_id = item["self_ref"]
            caption = ""
            for ann in item.get("annotations", []):
                if ann.get("kind") == "caption":
                    caption = ann.get("text", "")
                    break
            img_path = _save_picture_from_json(item, tmp_dir, element_id)
            pic_info.append({
                "element_id": element_id,
                "image_path": img_path,
                "page": page,
                "section": current_section,
                "caption": caption,
                "mime_type": "image/png",
            })
            elements.append({
                "type": "picture",
                "text": caption,
                "page": page,
                "section": current_section,
                "element_id": element_id,
            })

        elif label in ("list_item",):
            elements.append({"type": "list_item", "text": text, "page": page, "section": current_section})

        elif label in ("code",):
            elements.append({"type": "code", "text": text, "page": page, "section": current_section})

        elif label in ("caption", "footnote", "page_footer"):
            # Append to the preceding element if sensible
            if elements and elements[-1]["type"] in ("text", "table", "picture"):
                elements[-1]["text"] = (elements[-1]["text"] + "\n" + text).strip()

        elif text:
            elements.append({"type": "text", "text": text, "page": page, "section": current_section})

    return elements, pic_info


def _save_picture_from_json(pic_item: dict, tmp_dir: str, element_id: str) -> str | None:
    """Decode and save an embedded picture; return the file path or None."""
    data = pic_item.get("data") or {}
    uri = data.get("uri", "")
    if not uri.startswith("data:"):
        return None
    try:
        header, b64data = uri.split(",", 1)
        mime = header.split(";")[0].split(":")[1]
        ext = mimetypes.guess_extension(mime) or ".png"
        safe_id = element_id.replace("/", "_").replace("#", "")
        out_path = os.path.join(tmp_dir, f"{safe_id}{ext}")
        with open(out_path, "wb") as fh:
            fh.write(base64.b64decode(b64data))
        return out_path
    except Exception:
        return None


# ---------------------------------------------------------------------------
# Markdown fallback parser
# ---------------------------------------------------------------------------

def _parse_markdown_text(md: str) -> list[dict]:
    """Minimal markdown → elements without docling JSON."""
    import re
    elements: list[dict] = []
    current_section = ""
    for line in md.splitlines():
        stripped = line.rstrip()
        if re.match(r"^#{1,6}\s", stripped):
            current_section = stripped.lstrip("#").strip()
            elements.append({"type": "heading", "text": current_section, "page": 0, "section": current_section})
        elif stripped:
            elements.append({"type": "text", "text": stripped, "page": 0, "section": current_section})
    return elements


# ---------------------------------------------------------------------------
# Excel / CSV parsers
# ---------------------------------------------------------------------------

def _parse_excel(path: str) -> tuple[list, list]:
    if not _PANDAS_AVAILABLE:
        raise ImportError("pandas is required to parse Excel files")
    elements: list[dict] = []
    xf = pd.ExcelFile(path)
    for sheet_name in xf.sheet_names:
        df = xf.parse(sheet_name)
        table_data = _dataframe_to_table_data(df)
        elements.append({
            "type": "table",
            "text": "",
            "page": 0,
            "section": str(sheet_name),
            "table_data": table_data,
            "rows": table_data["num_rows"],
            "cols": table_data["num_cols"],
        })
    return elements, []


def _parse_csv(path: str) -> tuple[list, list]:
    rows: list[list[str]] = []
    with open(path, newline="", encoding="utf-8-sig") as fh:
        reader = csv.reader(fh)
        for row in reader:
            rows.append(row)
    if not rows:
        return [], []
    num_rows = len(rows)
    num_cols = max(len(r) for r in rows)
    cells: list[dict] = []
    for r_idx, row in enumerate(rows):
        for c_idx, cell_text in enumerate(row):
            cells.append({
                "start_row_offset_idx": r_idx,
                "end_row_offset_idx": r_idx + 1,
                "start_col_offset_idx": c_idx,
                "end_col_offset_idx": c_idx + 1,
                "text": cell_text,
                "column_header": r_idx == 0,
                "row_header": False,
            })
    table_data = {"num_rows": num_rows, "num_cols": num_cols, "table_cells": cells}
    stem = Path(path).stem
    return [{"type": "table", "text": "", "page": 0, "section": stem,
              "table_data": table_data, "rows": num_rows, "cols": num_cols}], []


def _dataframe_to_table_data(df) -> dict:
    """Convert a pandas DataFrame to docling-style table_data."""
    rows = [list(df.columns)] + df.astype(str).values.tolist()
    num_rows = len(rows)
    num_cols = max(len(r) for r in rows) if rows else 0
    cells: list[dict] = []
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cells.append({
                "start_row_offset_idx": r_idx,
                "end_row_offset_idx": r_idx + 1,
                "start_col_offset_idx": c_idx,
                "end_col_offset_idx": c_idx + 1,
                "text": str(val),
                "column_header": r_idx == 0,
                "row_header": False,
            })
    return {"num_rows": num_rows, "num_cols": num_cols, "table_cells": cells}


# ---------------------------------------------------------------------------
# PPTX parser
# ---------------------------------------------------------------------------

def _parse_pptx(path: str, tmp_dir: str) -> tuple[list, list]:
    if not _PPTX_AVAILABLE:
        raise ImportError("python-pptx is required to parse PPTX files")
    prs = Presentation(path)
    elements: list[dict] = []
    pic_info: list[dict] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        section = f"Slide {slide_num}"
        elements.append({"type": "heading", "text": section, "page": slide_num, "section": section})

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        elements.append({"type": "text", "text": text, "page": slide_num, "section": section})
            if shape.has_table:
                tbl = shape.table
                rows_data = [[cell.text for cell in row.cells] for row in tbl.rows]
                table_data = _rows_to_table_data(rows_data)
                elements.append({
                    "type": "table",
                    "text": "",
                    "page": slide_num,
                    "section": section,
                    "table_data": table_data,
                    "rows": table_data["num_rows"],
                    "cols": table_data["num_cols"],
                })

    return elements, pic_info


def _rows_to_table_data(rows: list[list[str]]) -> dict:
    num_rows = len(rows)
    num_cols = max(len(r) for r in rows) if rows else 0
    cells: list[dict] = []
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cells.append({
                "start_row_offset_idx": r_idx,
                "end_row_offset_idx": r_idx + 1,
                "start_col_offset_idx": c_idx,
                "end_col_offset_idx": c_idx + 1,
                "text": str(val),
                "column_header": r_idx == 0,
                "row_header": False,
            })
    return {"num_rows": num_rows, "num_cols": num_cols, "table_cells": cells}


# ---------------------------------------------------------------------------
# Image parser
# ---------------------------------------------------------------------------

def _parse_image(path: str, tmp_dir: str) -> tuple[list, list]:
    """Parse an image file into a picture element and pic_info."""
    # Determine mime type
    mime_type, _ = mimetypes.guess_type(path)
    if mime_type is None:
        mime_type = "image/png"

    # Read image file
    with open(path, "rb") as fh:
        image_data = fh.read()

    # Save to tmp_dir
    safe_name = Path(path).name.replace(" ", "_").replace("/", "_").replace("#", "")
    out_path = os.path.join(tmp_dir, safe_name)
    with open(out_path, "wb") as fh:
        fh.write(image_data)

    # Create element and pic_info
    element_id = f"picture_{Path(path).stem}"

    pic_info = [{
        "element_id": element_id,
        "image_path": out_path,
        "page": 0,
        "section": "Image",
        "caption": "",
        "mime_type": mime_type,
    }]

    elements = [{
        "type": "picture",
        "text": "",
        "page": 0,
        "section": "Image",
        "element_id": element_id,
    }]

    return elements, pic_info


# ---------------------------------------------------------------------------
# JSON parser
# ---------------------------------------------------------------------------

def _parse_json(path: str) -> tuple[list, list]:
    with open(path, "r", encoding="utf-8") as fh:
        data = json.load(fh)
    text = json.dumps(data, ensure_ascii=False, indent=2)
    stem = Path(path).stem
    return [{"type": "text", "text": text, "page": 0, "section": stem}], []


# ---------------------------------------------------------------------------
# Plain text fallback
# ---------------------------------------------------------------------------

def _parse_plaintext(path: str) -> tuple[list, list]:
    with open(path, "r", encoding="utf-8", errors="replace") as fh:
        text = fh.read()
    stem = Path(path).stem
    return [{"type": "text", "text": text, "page": 0, "section": stem}], []
